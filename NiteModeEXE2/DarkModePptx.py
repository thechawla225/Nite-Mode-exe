from pptx import Presentation
from pptx.slide import Slides
from pptx.dml.color import RGBColor
import pptx.text.text
import os


def convert_pptx(filename,filePathName,outputpath):    
    filePathName = os.path.join(filePathName,filename)
    prs = Presentation(filePathName)


    for slide in prs.slides:
        '''Changing background color'''
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        '''Adding Color to the text'''
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            try:
                i = 0
                while(text_frame.paragraphs[i]):
                    p = text_frame.paragraphs[i]
                    p.font.color.rgb = RGBColor(255,255,255)
                    i = i+1
            except:
                pass


    '''Saving the File'''
    path_to_file = os.path.join(outputpath,'DarkFile.pptx')
    prs.save(path_to_file)